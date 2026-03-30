#!/usr/bin/env python3
"""Generate AI Presentation PPTX for Tronc Commun Scientifique - with images."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BG_DARK = RGBColor(0x0f, 0x0c, 0x29)
BG_MID = RGBColor(0x1a, 0x1a, 0x2e)
PURPLE = RGBColor(0x81, 0x8c, 0xf8)
LIGHT_PURPLE = RGBColor(0xc0, 0x84, 0xfc)
PINK = RGBColor(0xf4, 0x72, 0xb6)
WHITE = RGBColor(0xff, 0xff, 0xff)
WHITE_70 = RGBColor(0xb3, 0xb3, 0xb3)
WHITE_50 = RGBColor(0x80, 0x80, 0x80)
GREEN = RGBColor(0x34, 0xd3, 0x99)
RED = RGBColor(0xf8, 0x71, 0x71)
CARD_BG = RGBColor(0x1e, 0x1b, 0x4b)

IMG_DIR = 'images'


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(6)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_image(slide, img_name, left, top, width=None, height=None):
    path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)


# =================== SLIDE 1: TITLE ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

# Background image (robot/AI themed)
add_image(slide, 'ai-brain.jpg', Inches(8.5), Inches(1.5), width=Inches(4.5))

# Semi-transparent overlay on left
add_shape_bg(slide, Inches(0), Inches(0), Inches(8.5), Inches(7.5), BG_DARK)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(1.5),
             "L'Intelligence Artificielle", font_size=48, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(3.0), Inches(7), Inches(0.8),
             "Projet Scolaire - Informatique", font_size=24, color=WHITE_70)

# Student info
box = add_shape_bg(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.0), CARD_BG, 0.05)
add_text_box(slide, Inches(1.1), Inches(4.5), Inches(5), Inches(0.5),
             "Realise par", font_size=16, color=WHITE_70)
add_text_box(slide, Inches(1.1), Inches(4.9), Inches(5), Inches(0.6),
             "Ayoub Assouar", font_size=30, color=LIGHT_PURPLE, bold=True)
add_text_box(slide, Inches(1.1), Inches(5.5), Inches(5), Inches(0.5),
             "Tronc Commun Scientifique", font_size=16, color=WHITE_50)

add_shape_bg(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), LIGHT_PURPLE)


# =================== SLIDE 2: SOMMAIRE ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_MID)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
             "Sommaire", font_size=40, color=PURPLE, bold=True)

chapters = [
    "Qu'est-ce que l'IA ?",
    "Histoire de l'IA",
    "Les types d'IA",
    "Comment fonctionne l'IA ?",
    "Applications de l'IA",
    "L'IA dans la vie quotidienne",
    "Avantages et Inconvenients",
    "L'avenir de l'IA",
]

for i, ch in enumerate(chapters):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.8) + row * Inches(1.25)
    box = add_shape_bg(slide, x, y, Inches(5.8), Inches(1.0), CARD_BG, 0.03)
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.6))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = PURPLE
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text_box(slide, x + Inches(1.0), y + Inches(0.2), Inches(4.5), Inches(0.6),
                 ch, font_size=18, color=WHITE_70)


# =================== SLIDE 3: Qu'est-ce que l'IA ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 1", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "Qu'est-ce que l'Intelligence Artificielle ?", font_size=36, color=WHITE, bold=True)

# Text on left
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.5), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "L'intelligence artificielle (IA) est une branche de l'informatique qui vise a creer des systemes capables d'effectuer des taches qui necessitent normalement l'intelligence humaine."
p.font.size = Pt(16)
p.font.color.rgb = WHITE_70
p.font.name = 'Calibri'
p.space_after = Pt(12)

add_paragraph(tf, "Ces taches incluent :", font_size=16, color=WHITE_70)
for item in ["La reconnaissance vocale et d'images", "La prise de decision autonome", "La traduction de langues", "La perception visuelle", "L'apprentissage a partir de donnees"]:
    add_paragraph(tf, "  *  " + item, font_size=14, color=WHITE_70, space_before=Pt(2), space_after=Pt(2))

add_paragraph(tf, "", font_size=8, color=WHITE_70, space_before=Pt(4), space_after=Pt(4))
add_paragraph(tf, 'Le terme "Intelligence Artificielle" a ete invente par John McCarthy en 1956.', font_size=16, color=LIGHT_PURPLE, bold=True)

# Image on right
add_image(slide, 'ai-robot.jpg', Inches(8.0), Inches(2.0), width=Inches(4.8))


# =================== SLIDE 4: Histoire ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_MID)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 2", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "L'Histoire de l'Intelligence Artificielle", font_size=36, color=WHITE, bold=True)

# Chess image on right
add_image(slide, 'ai-chess.jpg', Inches(8.8), Inches(2.0), width=Inches(4.0))

timeline = [
    ("1950", "Alan Turing propose le 'Test de Turing'"),
    ("1956", "Conference de Dartmouth : naissance de l'IA"),
    ("1966", "ELIZA, le premier chatbot (MIT)"),
    ("1997", "Deep Blue bat Garry Kasparov aux echecs"),
    ("2011", "IBM Watson gagne a Jeopardy!"),
    ("2016", "AlphaGo bat le champion mondial de Go"),
    ("2022", "ChatGPT revolutionne l'IA generative"),
]

add_shape_bg(slide, Inches(1.5), Inches(1.9), Inches(0.06), Inches(5.2), PURPLE)

for i, (year, desc) in enumerate(timeline):
    y = Inches(1.9) + i * Inches(0.74)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.35), y, Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = LIGHT_PURPLE
    dot.line.fill.background()
    add_text_box(slide, Inches(2.1), y - Inches(0.02), Inches(1.2), Inches(0.4),
                 year, font_size=16, color=LIGHT_PURPLE, bold=True)
    add_text_box(slide, Inches(3.3), y - Inches(0.02), Inches(5.2), Inches(0.4),
                 desc, font_size=15, color=WHITE_70)


# =================== SLIDE 5: Types d'IA ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 3", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "Les Types d'Intelligence Artificielle", font_size=36, color=WHITE, bold=True)

types_data = [
    ("IA Faible (Narrow AI)", "Concue pour UNE tache specifique.\n\nExemples concrets :\n* Siri et Alexa\n* Filtres spam Gmail\n* Recommandations Netflix\n* Google Translate", PURPLE, "Le type le plus courant"),
    ("IA Forte (General AI)", "Pourrait comprendre et apprendre N'IMPORTE quelle tache comme un humain.\n\nExemples hypothetiques :\n* Robot qui apprend tout seul\n* IA qui ressent des emotions\n\nN'existe PAS encore.", LIGHT_PURPLE, "En cours de recherche"),
    ("Super IA (Super AI)", "Depasserait l'intelligence humaine dans TOUS les domaines.\n\nConcept theorique qui souleve des questions ethiques majeures.\n\nEstimation : 2050-2100 ?", PINK, "Concept theorique"),
]

for i, (title, desc, color, subtitle) in enumerate(types_data):
    x = Inches(0.5) + i * Inches(4.2)
    box = add_shape_bg(slide, x, Inches(2.0), Inches(3.9), Inches(5.0), CARD_BG, 0.04)
    icon_box = add_shape_bg(slide, x + Inches(0.3), Inches(2.3), Inches(3.3), Inches(0.5), color, 0.1)
    add_text_box(slide, x + Inches(0.3), Inches(2.3), Inches(3.3), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.0), Inches(3.3), Inches(3.5),
                 desc, font_size=12, color=WHITE_70)
    add_text_box(slide, x + Inches(0.3), Inches(6.5), Inches(3.3), Inches(0.4),
                 subtitle, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)


# =================== SLIDE 6: Comment ca marche ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_MID)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 4", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "Comment Fonctionne l'IA ?", font_size=36, color=WHITE, bold=True)

# Network image
add_image(slide, 'ai-network.jpg', Inches(7.5), Inches(1.8), width=Inches(5.5))

how_items = [
    ("1. Les Donnees", "L'IA a besoin de GRANDES quantites de donnees pour apprendre.\nExemple : Des millions de photos de chats pour apprendre a reconnaitre un chat."),
    ("2. Les Algorithmes", "Des regles mathematiques analysent les donnees.\nExemple : Un algorithme de tri qui classe les emails en spam ou non-spam."),
    ("3. Machine Learning", "L'IA s'ameliore en s'entrainant. Elle apprend de ses erreurs.\nExemple : YouTube apprend vos gouts en analysant ce que vous regardez."),
    ("4. Reseaux de Neurones", "Inspires du cerveau humain, avec des couches de neurones artificiels.\nExemple : La reconnaissance faciale de votre telephone."),
]

for i, (title, desc) in enumerate(how_items):
    y = Inches(2.0) + i * Inches(1.3)
    box = add_shape_bg(slide, Inches(0.6), y, Inches(6.5), Inches(1.15), CARD_BG, 0.03)
    add_text_box(slide, Inches(0.9), y + Inches(0.1), Inches(6), Inches(0.4),
                 title, font_size=17, color=LIGHT_PURPLE, bold=True)
    add_text_box(slide, Inches(0.9), y + Inches(0.5), Inches(6), Inches(0.6),
                 desc, font_size=12, color=WHITE_70)


# =================== SLIDE 7: Applications ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 5", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "Les Applications de l'IA", font_size=36, color=WHITE, bold=True)

# Medical image
add_image(slide, 'ai-medical.jpg', Inches(8.5), Inches(1.8), width=Inches(4.5))

apps = [
    ("Medecine", "Diagnostic par IA, analyse de radiographies, robots chirurgicaux (Da Vinci)"),
    ("Transport", "Tesla Autopilot, Google Waymo, GPS intelligents, Uber/Lyft"),
    ("Education", "ChatGPT pour les devoirs, Duolingo, Khan Academy, correction auto"),
    ("Jeux Video", "PNJ intelligents dans GTA/FIFA, generation de mondes (Minecraft IA)"),
    ("Finance", "Detection de fraude bancaire, trading automatique, chatbots"),
    ("Environnement", "Prevision meteo, detection incendies, agriculture par drones"),
]

for i, (title, desc) in enumerate(apps):
    col = 0
    y = Inches(2.0) + i * Inches(0.85)
    box = add_shape_bg(slide, Inches(0.5), y, Inches(7.5), Inches(0.72), CARD_BG, 0.02)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.2), Inches(0.32), Inches(0.32))
    dot.fill.solid()
    dot.fill.fore_color.rgb = PURPLE
    dot.line.fill.background()
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(2), Inches(0.5),
                 title, font_size=15, color=LIGHT_PURPLE, bold=True)
    add_text_box(slide, Inches(3.3), y + Inches(0.15), Inches(4.5), Inches(0.45),
                 desc, font_size=12, color=WHITE_70)


# =================== SLIDE 8: Vie quotidienne ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_MID)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 6", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "L'IA dans Notre Vie Quotidienne", font_size=36, color=WHITE, bold=True)

# Phone image on right
add_image(slide, 'ai-phone.jpg', Inches(8.8), Inches(2.0), width=Inches(4.0))

daily = [
    ("Assistants vocaux", "Siri, Google Assistant, Alexa - posez une question, l'IA repond"),
    ("Reseaux sociaux", "TikTok, Instagram, YouTube - l'IA choisit votre contenu"),
    ("Streaming", "Netflix, Spotify - recommandations personnalisees"),
    ("Photo/Video", "Mode portrait, filtres Snapchat, retouche automatique"),
    ("Traduction", "Google Translate - 100+ langues en temps reel"),
    ("IA Generative", "ChatGPT, DALL-E, Midjourney - creation de texte et images"),
]

for i, (title, desc) in enumerate(daily):
    y = Inches(2.0) + i * Inches(0.85)
    box = add_shape_bg(slide, Inches(0.6), y, Inches(7.8), Inches(0.72), CARD_BG, 0.02)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.22), Inches(0.28), Inches(0.28))
    dot.fill.solid()
    dot.fill.fore_color.rgb = PURPLE
    dot.line.fill.background()
    add_text_box(slide, Inches(1.4), y + Inches(0.12), Inches(2.5), Inches(0.5),
                 title, font_size=15, color=LIGHT_PURPLE, bold=True)
    add_text_box(slide, Inches(3.9), y + Inches(0.15), Inches(4.3), Inches(0.45),
                 desc, font_size=13, color=WHITE_70)


# =================== SLIDE 9: Avantages & Inconvenients ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 7", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "Avantages et Inconvenients de l'IA", font_size=36, color=WHITE, bold=True)

# Advantages
add_shape_bg(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0), CARD_BG, 0.04)
add_shape_bg(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(0.6), GREEN, 0.04)
add_text_box(slide, Inches(0.8), Inches(2.05), Inches(5.2), Inches(0.5),
             "AVANTAGES", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

pros = [
    "Automatisation des taches repetitives",
    "Analyse de donnees ultra-rapide",
    "Disponible 24h/24, jamais fatigue",
    "Diagnostics medicaux plus precis",
    "Reduction des erreurs humaines",
    "Innovation scientifique acceleree",
]
for i, pro in enumerate(pros):
    y = Inches(2.8) + i * Inches(0.65)
    add_text_box(slide, Inches(0.9), y, Inches(5.2), Inches(0.5),
                 "+ " + pro, font_size=14, color=WHITE_70)

# Disadvantages
add_shape_bg(slide, Inches(7.0), Inches(2.0), Inches(5.8), Inches(5.0), CARD_BG, 0.04)
add_shape_bg(slide, Inches(7.0), Inches(2.0), Inches(5.8), Inches(0.6), RED, 0.04)
add_text_box(slide, Inches(7.3), Inches(2.05), Inches(5.2), Inches(0.5),
             "INCONVENIENTS", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

cons = [
    "Perte d'emplois (chomage technologique)",
    "Vie privee menacee (surveillance)",
    "Biais et discrimination algorithmique",
    "Dependance excessive a la technologie",
    "Cout eleve de developpement",
    "Deepfakes et desinformation",
]
for i, con in enumerate(cons):
    y = Inches(2.8) + i * Inches(0.65)
    add_text_box(slide, Inches(7.4), y, Inches(5.2), Inches(0.5),
                 "- " + con, font_size=14, color=WHITE_70)


# =================== SLIDE 10: L'avenir ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_MID)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CHAPITRE 8", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
             "L'Avenir de l'Intelligence Artificielle", font_size=36, color=WHITE, bold=True)

# Future image
add_image(slide, 'ai-future.jpg', Inches(8.5), Inches(1.8), width=Inches(4.5))

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(7), Inches(0.5),
             "Tendances pour les prochaines annees :", font_size=17, color=WHITE_70)

future = [
    ("IA Generative", "Textes, images, videos et musique crees par IA - de plus en plus realistes"),
    ("Medecine personnalisee", "Traitements uniques pour chaque patient grace a l'analyse ADN par IA"),
    ("Robots intelligents", "Robots qui interagissent naturellement (Boston Dynamics, Tesla Optimus)"),
    ("IA ethique", "Nouvelles lois et reglementations pour encadrer l'utilisation de l'IA"),
    ("Education", "Professeurs IA personnels, adaptation au niveau de chaque eleve"),
]

for i, (title, desc) in enumerate(future):
    y = Inches(2.5) + i * Inches(0.95)
    box = add_shape_bg(slide, Inches(0.6), y, Inches(7.5), Inches(0.8), CARD_BG, 0.03)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.25), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = PINK
    dot.line.fill.background()
    add_text_box(slide, Inches(1.5), y + Inches(0.08), Inches(2.5), Inches(0.4),
                 title, font_size=15, color=LIGHT_PURPLE, bold=True)
    add_text_box(slide, Inches(1.5), y + Inches(0.4), Inches(6.4), Inches(0.35),
                 desc, font_size=12, color=WHITE_70)


# =================== SLIDE 11: Conclusion ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             "CONCLUSION", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(1.5), Inches(1.0), Inches(10.3), Inches(1.0),
             "L'IA : Un Outil Puissant a Utiliser avec Responsabilite",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.2),
             "L'intelligence artificielle transforme notre monde. Elle offre des possibilites immenses pour la medecine, l'education, l'environnement et bien plus. Mais elle souleve aussi des defis importants : vie privee, emploi, ethique.",
             font_size=17, color=WHITE_70, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(1.0),
             "En tant qu'etudiants et futurs citoyens, nous devons comprendre l'IA, connaitre ses avantages et ses limites, et contribuer a un developpement ethique et responsable de cette technologie.",
             font_size=17, color=WHITE_70, alignment=PP_ALIGN.CENTER)

# Quote
quote_box = add_shape_bg(slide, Inches(2.5), Inches(5.2), Inches(8.3), Inches(1.5), CARD_BG, 0.04)
add_text_box(slide, Inches(2.8), Inches(5.4), Inches(7.7), Inches(0.7),
             '"L\'intelligence artificielle est la nouvelle electricite."',
             font_size=22, color=LIGHT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.8), Inches(6.1), Inches(7.7), Inches(0.4),
             "- Andrew Ng, professeur a Stanford et expert mondial en IA",
             font_size=14, color=WHITE_50, alignment=PP_ALIGN.CENTER)


# =================== SLIDE 12: Merci ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.5),
             "Merci !", font_size=64, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.8),
             "Presentation realisee par Ayoub Assouar", font_size=22, color=WHITE_70, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.6),
             "Tronc Commun Scientifique - 2024/2025", font_size=18, color=WHITE_50, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.6),
             "Des questions ?", font_size=20, color=WHITE_50, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), GREEN)

# Save
prs.save('AI_Presentation_Ayoub_Assouar.pptx')
print("Presentation saved successfully: AI_Presentation_Ayoub_Assouar.pptx")
